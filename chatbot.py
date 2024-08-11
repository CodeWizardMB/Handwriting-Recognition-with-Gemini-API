import cv2
import google.generativeai as genai
import language_tool_python
import pyttsx3
import speech_recognition as sr
import os
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
from googletrans import Translator

# Get your API key from the environment variable
api_key = os.getenv("API_KEY")
if api_key is None:
    raise ValueError("API_KEY environment variable not set")

genai.configure(api_key=api_key)

def capture_image(image_path):
    """Captures an image from the webcam and saves it to the specified path."""
    cap = cv2.VideoCapture(0)

    if not cap.isOpened():
        print("Error: Could not open webcam.")
        return

    print("Press 's' to capture the image.")
    while True:
        ret, frame = cap.read()
        if not ret:
            print("Error: Could not read frame.")
            break

        cv2.imshow('Capture', frame)

        key = cv2.waitKey(1)
        if key == ord('s'):
            cv2.imwrite(image_path, frame)
            print(f"Image saved to {image_path}")
            break

    cap.release()
    cv2.destroyAllWindows()

def upload_to_gemini(path, mime_type=None):
    """Uploads the given file to Gemini."""
    file = genai.upload_file(path, mime_type=mime_type)
    print(f"Uploaded file '{file.display_name}' as: {file.uri}")
    return file

def read_text_from_document(file_path):
    """Reads text from .txt, .docx, .pptx, or .pdf file."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.txt':
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    elif ext == '.docx':
        doc = Document(file_path)
        return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
    elif ext == '.pptx':
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    elif ext == '.pdf':
        doc = fitz.open(file_path)
        text = []
        for page in doc:
            text.append(page.get_text())
        return '\n'.join(text)
    else:
        print("Unsupported file type.")
        return ""

def correct_grammar_with_languagetool(text):
    """Corrects grammar using LanguageTool."""
    tool = language_tool_python.LanguageTool('en-US')
    matches = tool.check(text)
    corrected_text = language_tool_python.utils.correct(text, matches)
    return corrected_text

def translate_text(text, target_language):
    """Translates text to the target language using Google Translate."""
    translator = Translator()
    translation = translator.translate(text, dest=target_language)
    return translation.text

def speak_text(text, rate=200):
    """Converts text to speech."""
    engine = pyttsx3.init()
    engine.setProperty('rate', rate)
    engine.say(text)
    engine.runAndWait()

def recognize_speech_from_microphone():
    """Recognizes speech from the microphone."""
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        print("Listening for command...")
        audio = recognizer.listen(source)
        try:
            command = recognizer.recognize_google(audio)
            print(f"Recognized command: {command}")
            return command.lower()
        except sr.UnknownValueError:
            print("Sorry, I did not understand the audio.")
            return ""
        except sr.RequestError:
            print("Sorry, there was an issue with the speech recognition service.")
            return ""

# Ask the user to select an option
print("Select an option:")
print("1. Capture a new image")
print("2. Upload an existing image or document")
choice = input("Enter 1 or 2: ")

if choice == '1':
    image_path = "captured_image.jpg"
    capture_image(image_path)
    file_to_upload = image_path
elif choice == '2':
    print("Select the document type to upload:")
    print("1. Image (.jpg, .png)")
    print("2. Text (.txt)")
    print("3. Word (.docx)")
    print("4. PowerPoint (.pptx)")
    print("5. PDF (.pdf)")
    doc_choice = input("Enter 1, 2, 3, 4, or 5: ")

    if doc_choice == '1':
        file_to_upload = input("Enter the path of the image to upload: ")
    elif doc_choice == '2':
        file_to_upload = input("Enter the path of the text file to upload: ")
    elif doc_choice == '3':
        file_to_upload = input("Enter the path of the Word document to upload: ")
    elif doc_choice == '4':
        file_to_upload = input("Enter the path of the PowerPoint presentation to upload: ")
    elif doc_choice == '5':
        file_to_upload = input("Enter the path of the PDF file to upload: ")
    else:
        print("Invalid choice.")
        exit()

# Check if the file exists
if os.path.exists(file_to_upload):
    mime_type = None
    ext = os.path.splitext(file_to_upload)[1].lower()
    if ext in ['.jpg', '.png']:
        mime_type = 'image/jpeg'
    elif ext == '.txt':
        mime_type = 'text/plain'
    elif ext == '.docx':
        mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    elif ext == '.pptx':
        mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    elif ext == '.pdf':
        mime_type = 'application/pdf'
    else:
        print("Unsupported file type.")
        exit()

    # Upload the file
    file = upload_to_gemini(file_to_upload, mime_type=mime_type)

    # Create the model
    generation_config = {
        "temperature": 1,
        "top_p": 0.95,
        "top_k": 64,
        "max_output_tokens": 8192,
        "response_mime_type": "text/plain",
    }

    model = genai.GenerativeModel(
        model_name="gemini-1.5-pro",
        generation_config=generation_config,
    )

    # Start chat session
    chat_session = model.start_chat(
        history=[
            {
                "role": "user",
                "parts": [
                    file.uri,
                    "Please recognize the text in this file.",
                ],
            },
        ]
    )

    # Continuously interact with the chatbot
    while True:
        user_input = input("You: ")
        if user_input.lower() in ["exit", "quit"]:
            print("Exiting chat.")
            break

        # Send user input as a message to the chat session
        response = chat_session.send_message(user_input)
        recognized_text = response.text
        print("Bot:", recognized_text)

        # Convert bot response to speech
        speak_text(recognized_text)

else:
    print("The specified file does not exist.")
